"""
CampusSpeak AI - Flask Application
Real-time speech analysis with Web Speech API + Python backend.
"""

import os
import re
import json
import logging
import requests as http_requests
from flask import Flask, render_template, request, jsonify, redirect, session, url_for
from flask_login import LoginManager, login_user, logout_user, login_required, current_user
from dotenv import load_dotenv
from models import db, User, Session, PracticeSession
from datetime import datetime, timedelta
from functools import wraps

load_dotenv()

app = Flask(__name__)
# Keep a stable dev fallback to avoid session invalidation on app reloads.
app.secret_key = os.getenv("SECRET_KEY", "dev-secret-key-change-in-production")

# ─── Database Configuration ──────────────────────────────

app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///presentation_coach.db'
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

db.init_app(app)

# ─── Login Manager Setup ─────────────────────────────────

login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'
login_manager.login_message = 'Please log in to access this page.'

@login_manager.user_loader
def load_user(user_id):
    """Load user by ID for Flask-Login."""
    return User.query.get(int(user_id))

# ─── AI Provider Setup ─────────────────────────────────────

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
GEMINI_MODEL = os.getenv("GEMINI_MODEL", "gemini-2.0-flash")
GEMINI_URL = "https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"

OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY", "")
OPENROUTER_MODEL = os.getenv("OPENROUTER_MODEL", "nvidia/nemotron-3-super-120b-a12b:free")
OPENROUTER_URL = "https://openrouter.ai/api/v1/chat/completions"
USE_OPENROUTER_FALLBACK = os.getenv("USE_OPENROUTER_FALLBACK", "false").lower() == "true"

COACH_SYSTEM_PROMPT = (
    "You are Zoiee, a friendly and encouraging CampusSpeak AI coach robot. "
    "You help users improve their public speaking and presentation skills. "
    "Keep responses concise (2-3 sentences max). Use a warm, supportive tone. "
    "When given speech metrics (WPM, clarity, filler words), provide specific, "
    "actionable advice. If the user asks something unrelated to presentations, "
    "gently steer back to coaching. Never use markdown formatting in responses."
)

FILLER_WORDS = [
    "um", "uh", "like", "you know", "basically", "actually",
    "literally", "so", "well", "right", "I mean", "kind of",
    "sort of", "you see",
]


def build_chat_messages(user_message: str, context: dict, history: list) -> list:
    """Build chat messages with prompt, metrics context, and recent history."""
    messages = [{"role": "system", "content": COACH_SYSTEM_PROMPT}]

    if context.get("isRecording"):
        metrics_info = (
            f"Current session metrics - WPM: {context.get('wpm', 0)}, "
            f"Clarity: {context.get('clarity', 0)}%, "
            f"Filler words: {context.get('fillers', 0)}"
        )
        messages.append({"role": "system", "content": metrics_info})

    for msg in history:
        role = "user" if msg.get("role") == "user" else "assistant"
        messages.append({"role": role, "content": str(msg.get("content", ""))[:500]})

    messages.append({"role": "user", "content": user_message})
    return messages


def _extract_gemini_text(result: dict) -> str:
    """Extract assistant text from Gemini generateContent response."""
    candidates = result.get("candidates") or []
    if not candidates:
        return ""
    content = candidates[0].get("content") or {}
    parts = content.get("parts") or []
    text_chunks = [p.get("text", "") for p in parts if isinstance(p, dict) and p.get("text")]
    return "\n".join(text_chunks).strip()


def request_gemini(messages: list) -> str:
    """Call Gemini API and return reply text; raise on failure."""
    system_parts = [m["content"] for m in messages if m.get("role") == "system"]
    conversation_parts = []
    for m in messages:
        if m.get("role") == "system":
            continue
        prefix = "User" if m.get("role") == "user" else "Coach"
        conversation_parts.append(f"{prefix}: {m.get('content', '')}")

    prompt_sections = []
    if system_parts:
        prompt_sections.append("System Instructions:\n" + "\n".join(system_parts))
    if conversation_parts:
        prompt_sections.append("Conversation:\n" + "\n".join(conversation_parts))
    prompt_sections.append("Respond as Coach in 2-3 concise sentences.")
    full_prompt = "\n\n".join(prompt_sections)

    response = http_requests.post(
        url=GEMINI_URL.format(model=GEMINI_MODEL) + f"?key={GEMINI_API_KEY}",
        headers={"Content-Type": "application/json"},
        data=json.dumps({
            "contents": [{"parts": [{"text": full_prompt}]}],
            "generationConfig": {
                "temperature": 0.7,
                "maxOutputTokens": 200,
            },
        }),
        timeout=30,
    )
    result = response.json()
    if response.status_code != 200:
        error_msg = result.get("error", {}).get("message", "Unknown error")
        raise RuntimeError(f"Gemini API error: {error_msg}")

    reply = _extract_gemini_text(result)
    if not reply:
        raise RuntimeError("Gemini API returned empty response")
    return reply


def request_openrouter(messages: list) -> str:
    """Call OpenRouter API and return reply text; raise on failure."""
    response = http_requests.post(
        url=OPENROUTER_URL,
        headers={
            "Authorization": f"Bearer {OPENROUTER_API_KEY}",
            "Content-Type": "application/json",
        },
        data=json.dumps({
            "model": OPENROUTER_MODEL,
            "messages": messages,
        }),
        timeout=30,
    )
    result = response.json()

    if response.status_code != 200:
        error_msg = result.get("error", {}).get("message", "Unknown error")
        raise RuntimeError(f"OpenRouter API error: {error_msg}")

    choices = result.get("choices") or []
    if not choices:
        raise RuntimeError("OpenRouter API returned no choices")
    message = choices[0].get("message") or {}
    reply = str(message.get("content", "")).strip()
    if not reply:
        raise RuntimeError("OpenRouter API returned empty response")
    return reply


def analyze_speech(text: str, duration_sec: float) -> dict:
    """Analyze speech text and return detailed feedback metrics."""
    words = text.strip().split()
    word_count = len(words)
    
    wpm = 0
    if duration_sec > 3:
        wpm = round((word_count / duration_sec) * 60)
    elif duration_sec > 0:
        # Prevent extreme WPM spikes during the first 3 seconds
        wpm = min(150, round((word_count / max(1.0, duration_sec)) * 60))

    # Count filler words
    filler_count = 0
    filler_details = {}
    lower_text = text.lower()
    for filler in FILLER_WORDS:
        matches = re.findall(rf"\b{re.escape(filler)}\b", lower_text)
        if matches:
            filler_details[filler] = len(matches)
            filler_count += len(matches)

    # Clarity score (heuristic-based)
    clarity = 100
    if word_count > 0:
        filler_ratio = filler_count / word_count
        clarity = max(0, round(100 - filler_ratio * 500))
    if wpm > 180:
        clarity = max(0, clarity - 15)
    if 0 < wpm < 100:
        clarity = max(0, clarity - 10)

    # Pace rating
    if wpm == 0:
        pace = "N/A"
    elif wpm < 100:
        pace = "Too Slow"
    elif wpm <= 150:
        pace = "Good"
    elif wpm <= 170:
        pace = "Slightly Fast"
    else:
        pace = "Too Fast"

    # ── Confidence Score ───────────────────────────────────────────────────
    # Composite score: pace contributes 40%, clarity 40%, filler penalty 20%
    pace_score = 0
    if wpm >= 120 and wpm <= 150:
        pace_score = 100
    elif wpm >= 100 and wpm < 120:
        pace_score = 75
    elif wpm > 150 and wpm <= 170:
        pace_score = 75
    elif wpm > 170:
        pace_score = 50
    elif wpm > 0:
        pace_score = 40

    filler_penalty = min(100, filler_count * 10)
    filler_score = max(0, 100 - filler_penalty)

    if word_count == 0:
        confidence_score = 0
    else:
        confidence_score = round(pace_score * 0.4 + clarity * 0.4 + filler_score * 0.2)
        confidence_score = max(0, min(100, confidence_score))

    # ── Personality / Communication Type ──────────────────────────────────
    if word_count == 0:
        personality_type = "Unknown"
        improvement_area = "Start practicing to get your personality analysis!"
    elif filler_count > 6 and (wpm < 110 or wpm == 0):
        personality_type = "Hesitant Speaker"
        improvement_area = "Confidence & Fluency — replace fillers with deliberate pauses"
    elif wpm > 170:
        personality_type = "Enthusiastic Speaker"
        improvement_area = "Pacing & Clarity — slow down to let ideas land"
    elif filler_count > 4:
        personality_type = "Nervous Speaker"
        improvement_area = "Filler Word Reduction — practice mindful pausing"
    elif wpm >= 120 and wpm <= 150 and filler_count <= 2:
        personality_type = "Confident Communicator"
        improvement_area = "Storytelling & Intonation — add more vocal variety"
    elif wpm < 100 and wpm > 0:
        personality_type = "Thoughtful Speaker"
        improvement_area = "Engagement & Energy — increase pace slightly"
    elif clarity >= 85:
        personality_type = "Analytical Speaker"
        improvement_area = "Connection & Eye Contact — engage your audience more"
    else:
        personality_type = "Developing Speaker"
        improvement_area = "Overall Fluency — practice daily with short sessions"

    # Suggestions
    suggestions = []
    if wpm > 170:
        suggestions.append(
            "You're speaking quite fast. Try slowing down to ~140 WPM "
            "for better audience comprehension."
        )
    if 0 < wpm < 100:
        suggestions.append(
            "Your pace is a bit slow. Aim for 120–150 WPM to keep "
            "the audience engaged."
        )
    if filler_count > 5:
        suggestions.append(
            f"You used {filler_count} filler words. Practice replacing "
            f"them with a brief pause."
        )
    elif filler_count > 2:
        suggestions.append(
            f"You used {filler_count} filler words — try to reduce them."
        )
    if filler_count == 0 and word_count > 10:
        suggestions.append("Great job avoiding filler words! Keep it up.")
    if duration_sec < 30 and word_count > 0:
        suggestions.append(
            "Your session was very short. Practice for at least 2 minutes "
            "for meaningful feedback."
        )
    if clarity >= 90 and word_count > 10:
        suggestions.append("Excellent clarity! Your speech is clean and well-paced.")
    if not suggestions:
        suggestions.append("Keep practicing to improve your presentation skills!")

    return {
        "word_count": word_count,
        "wpm": wpm,
        "filler_count": filler_count,
        "filler_details": filler_details,
        "clarity": clarity,
        "pace": pace,
        "duration": round(duration_sec, 1),
        "suggestions": suggestions,
        "confidence_score": confidence_score,
        "personality_type": personality_type,
        "improvement_area": improvement_area,
    }


# ─── Authentication Routes ────────────────────────────

@app.route("/")
def home():
    """Redirect to login page - login is primary entry point."""
    return redirect("/login")


@app.route("/login", methods=['GET', 'POST'])
def login():
    """Login page and authentication endpoint."""
    if request.method == 'POST':
        data = request.get_json()
        email = data.get('email', '').lower().strip()
        password = data.get('password', '')
        remember_me = data.get('remember_me', False)
        
        if not email or not password:
            return jsonify({'error': 'Email and password required'}), 400
        
        user = User.query.filter_by(email=email).first()
        
        if user and user.check_password(password):
            # Update last login
            user.update_last_login()
            
            # Handle remember me
            if remember_me:
                remember_token = user.generate_remember_token()
                db.session.commit()
                session['remember_token'] = remember_token
            
            # Log user in
            login_user(user, remember=remember_me)
            db.session.commit()
            
            return jsonify({'success': True, 'message': 'Login successful'}), 200
        else:
            return jsonify({'error': 'Invalid email or password'}), 401
    
    return render_template("login.html")


@app.route("/register", methods=['POST'])
def register():
    """User registration endpoint."""
    data = request.get_json()
    
    email = data.get('email', '').lower().strip()
    username = data.get('username', '').strip()
    password = data.get('password', '')
    full_name = data.get('full_name', '').strip()
    
    # Validation
    if not all([email, username, password, full_name]):
        return jsonify({'error': 'All fields required'}), 400
    
    if len(password) < 8:
        return jsonify({'error': 'Password must be at least 8 characters'}), 400
    
    if User.query.filter_by(email=email).first():
        return jsonify({'error': 'Email already registered'}), 409
    
    if User.query.filter_by(username=username).first():
        return jsonify({'error': 'Username already taken'}), 409
    
    # Create new user
    user = User(
        email=email,
        username=username,
        full_name=full_name
    )
    user.set_password(password)
    user.generate_verification_token()
    
    db.session.add(user)
    db.session.commit()
    
    return jsonify({
        'success': True,
        'message': 'Registration successful! Please verify your email.'
    }), 201


@app.route("/verify-email", methods=['POST'])
def verify_email():
    """Verify email with token."""
    data = request.get_json()
    token = data.get('token', '')
    
    user = User.query.filter_by(verification_token=token).first()
    
    if not user:
        return jsonify({'error': 'Invalid verification token'}), 400
    
    if user.verification_token_expires < datetime.utcnow():
        return jsonify({'error': 'Verification token expired'}), 400
    
    user.verify_email()
    db.session.commit()
    
    return jsonify({'success': True, 'message': 'Email verified successfully'}), 200


@app.route("/forgot-password", methods=['POST'])
def forgot_password():
    """Request password reset."""
    data = request.get_json()
    email = data.get('email', '').lower().strip()
    
    user = User.query.filter_by(email=email).first()
    
    if user:
        reset_token = user.generate_reset_token()
        db.session.commit()
        
        # TODO: Send reset email with token
        # For demo, we'll return the token
        return jsonify({
            'success': True,
            'message': 'Password reset link sent to email',
            'reset_token': reset_token  # Remove in production!
        }), 200
    
    # Security: Don't reveal if email exists
    return jsonify({
        'success': True,
        'message': 'Password reset link sent to email'
    }), 200


@app.route("/reset-password", methods=['POST'])
def reset_password():
    """Reset password with token."""
    data = request.get_json()
    token = data.get('token', '')
    new_password = data.get('new_password', '')
    
    if len(new_password) < 8:
        return jsonify({'error': 'Password must be at least 8 characters'}), 400
    
    user = User.query.filter_by(reset_token=token).first()
    
    if not user:
        return jsonify({'error': 'Invalid reset token'}), 400
    
    if not user.verify_reset_token(token):
        return jsonify({'error': 'Reset token expired'}), 400
    
    user.reset_password(new_password)
    db.session.commit()
    
    return jsonify({'success': True, 'message': 'Password reset successful'}), 200


@app.route("/logout", methods=['GET', 'POST'])
def logout():
    """Logout user."""
    logout_user()
    session.clear()
    return redirect("/login")


@app.route("/profile")
@login_required
def profile():
    """User profile page."""
    return jsonify(current_user.to_dict()), 200


@app.route("/app")
@login_required
def app_dashboard():
    """Main presentation coach application - requires login."""
    return render_template("index.html")


@app.route("/demo")
def demo():
    """Demo mode - app without authentication required."""
    return render_template("index.html", demo_mode=True)


@app.route("/dashboard")
def dashboard():
    """Progress dashboard.

    - Authenticated users see their real dashboard.
    - Demo mode users can access a demo dashboard via /dashboard?demo=1.
    """
    is_demo = request.args.get("demo") == "1"

    if current_user.is_authenticated:
        return render_template("dashboard.html", demo_mode=False)

    if is_demo:
        return render_template("dashboard.html", demo_mode=True)

    return redirect(url_for("login"))


@app.route("/save-session", methods=["POST"])
@login_required
def save_session():
    """Save a completed practice session's metrics to the database."""
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "Invalid JSON body"}), 400

    try:
        ps = PracticeSession(
            user_id=current_user.id,
            duration=float(data.get("duration", 0)),
            word_count=int(data.get("word_count", 0)),
            wpm=int(data.get("wpm", 0)),
            filler_count=int(data.get("filler_count", 0)),
            clarity=int(data.get("clarity", 0)),
            confidence_score=int(data.get("confidence_score", 0)),
            pace=str(data.get("pace", "N/A")),
            personality_type=str(data.get("personality_type", "")),
            improvement_area=str(data.get("improvement_area", "")),
            mode=str(data.get("mode", "Free Speaking")),
        )
        db.session.add(ps)
        db.session.commit()
        return jsonify({"success": True, "session_id": ps.id}), 201
    except Exception as e:
        db.session.rollback()
        logging.error("Failed to save session: %s", str(e))
        return jsonify({"error": "Failed to save session"}), 500


@app.route("/dashboard-data", methods=["GET"])
@login_required
def dashboard_data():
    """Return the last 30 practice sessions for the current user."""
    sessions = (
        PracticeSession.query
        .filter_by(user_id=current_user.id)
        .order_by(PracticeSession.created_at.asc())
        .limit(30)
        .all()
    )
    return jsonify({
        "username": current_user.full_name or current_user.username,
        "sessions": [s.to_dict() for s in sessions],
        "total_sessions": len(sessions),
    })


@app.route("/demo-dashboard-data", methods=["GET"])
def demo_dashboard_data():
    """Return sample dashboard data for demo mode (no login required)."""
    now = datetime.utcnow()
    sessions = [
        {
            "id": 1,
            "duration": 54.2,
            "word_count": 105,
            "wpm": 116,
            "filler_count": 7,
            "clarity": 72,
            "confidence_score": 64,
            "pace": "Slightly Fast",
            "personality_type": "Developing Speaker",
            "improvement_area": "Overall Fluency - practice daily with short sessions",
            "mode": "Self Introduction",
            "created_at": (now - timedelta(days=6)).isoformat() + "Z",
        },
        {
            "id": 2,
            "duration": 62.0,
            "word_count": 118,
            "wpm": 121,
            "filler_count": 6,
            "clarity": 76,
            "confidence_score": 69,
            "pace": "Good",
            "personality_type": "Thoughtful Speaker",
            "improvement_area": "Engagement & Energy - increase pace slightly",
            "mode": "HR Interview",
            "created_at": (now - timedelta(days=4)).isoformat() + "Z",
        },
        {
            "id": 3,
            "duration": 68.5,
            "word_count": 132,
            "wpm": 126,
            "filler_count": 5,
            "clarity": 80,
            "confidence_score": 74,
            "pace": "Good",
            "personality_type": "Analytical Speaker",
            "improvement_area": "Connection & Eye Contact - engage your audience more",
            "mode": "Presentation",
            "created_at": (now - timedelta(days=3)).isoformat() + "Z",
        },
        {
            "id": 4,
            "duration": 71.4,
            "word_count": 140,
            "wpm": 132,
            "filler_count": 4,
            "clarity": 84,
            "confidence_score": 79,
            "pace": "Good",
            "personality_type": "Confident Communicator",
            "improvement_area": "Storytelling & Intonation - add more vocal variety",
            "mode": "Group Discussion",
            "created_at": (now - timedelta(days=1)).isoformat() + "Z",
        },
        {
            "id": 5,
            "duration": 67.8,
            "word_count": 145,
            "wpm": 128,
            "filler_count": 4,
            "clarity": 82,
            "confidence_score": 82,
            "pace": "Good",
            "personality_type": "Confident Communicator",
            "improvement_area": "Storytelling & Intonation - add more vocal variety",
            "mode": "Presentation",
            "created_at": now.isoformat() + "Z",
        },
    ]

    return jsonify({
        "username": "Demo User",
        "sessions": sessions,
        "total_sessions": len(sessions),
        "is_demo": True,
    })


# ─── Routes ───────────────────────────────────────────────


ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.ppt', '.pptx'}
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10 MB


@app.route("/import-file", methods=["POST"])
def import_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    if not file.filename:
        return jsonify({"error": "Empty filename"}), 400

    ext = os.path.splitext(file.filename)[1].lower()
    if ext == '.doc':
        return jsonify({"error": "Legacy .doc files are not supported. Please convert to .docx first."}), 400

    if ext not in ALLOWED_EXTENSIONS:
        return jsonify({"error": f"Unsupported file type: {ext}"}), 400

    # Read file content with size check
    file_bytes = file.read(MAX_FILE_SIZE + 1)
    if len(file_bytes) > MAX_FILE_SIZE:
        return jsonify({"error": "File too large (max 10 MB)"}), 400

    try:
        if ext == '.pdf':
            import PyPDF2
            import io
            reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
        elif ext == '.docx':
            import docx
            import io
            doc = docx.Document(io.BytesIO(file_bytes))
            text = "\n".join(p.text for p in doc.paragraphs)
        elif ext in ('.ppt', '.pptx'):
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(file_bytes))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                parts = [f"--- Slide {i} ---"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
                slides_text.append("\n".join(parts))
            text = "\n\n".join(slides_text)
        else:
            return jsonify({"error": "Unsupported format"}), 400

        return jsonify({"text": text.strip()})
    except Exception as e:
        return jsonify({"error": f"Failed to parse file: {str(e)}"}), 500


@app.route("/analyze", methods=["POST"])
def analyze():
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "Invalid JSON body"}), 400

    text = str(data.get("text", ""))
    try:
        duration = float(data.get("duration", 0))
    except (TypeError, ValueError):
        return jsonify({"error": "Invalid duration value"}), 400

    result = analyze_speech(text, duration)
    return jsonify(result)


@app.route("/chat", methods=["POST"])
def chat():
    if not GEMINI_API_KEY and not OPENROUTER_API_KEY:
        return jsonify({"reply": "No AI provider configured. Please set GEMINI_API_KEY or OPENROUTER_API_KEY in your .env file."}), 200

    data = request.get_json(silent=True)
    if not data or not data.get("message", "").strip():
        return jsonify({"error": "Message is required"}), 400

    user_message = str(data["message"]).strip()[:1000]
    context = data.get("context", {})
    history = data.get("history", [])[-10:]

    messages = build_chat_messages(user_message, context, history)

    try:
        # Gemini is the default provider.
        if GEMINI_API_KEY:
            try:
                reply = request_gemini(messages)
                return jsonify({"reply": reply, "provider": "gemini"})
            except Exception as gemini_error:
                logging.error("Gemini request failed: %s", str(gemini_error))
                if not USE_OPENROUTER_FALLBACK:
                    return jsonify({"reply": "Gemini is configured but unavailable right now. Please try again in a moment.", "provider": "gemini"}), 200

        # OpenRouter fallback is optional and off by default.
        if USE_OPENROUTER_FALLBACK and OPENROUTER_API_KEY:
            try:
                reply = request_openrouter(messages)
                return jsonify({"reply": reply, "provider": "openrouter"})
            except Exception as openrouter_error:
                logging.error("OpenRouter fallback failed: %s", str(openrouter_error))

        return jsonify({"reply": "Sorry, I'm having trouble connecting right now. Please try again!"}), 200
    except Exception as e:
        logging.error("Chat request failed: %s", str(e))
        return jsonify({"reply": "Sorry, I'm having trouble connecting right now. Please try again!"}), 200


# ─── Demo Session Endpoints ───────────────────────────────

DEMO_SCRIPT = """Good morning, everyone! Today I want to talk about artificial intelligence and how it's transforming our world.
First, let me start with the basics. AI is not a new concept, but recent advancements have made it incredibly powerful.
There are three key applications we should consider: healthcare, education, and business automation.
In healthcare, AI is helping doctors diagnose diseases faster and more accurately.
In education, personalized learning systems are adapting to each student's needs.
And in business, automation is increasing efficiency and reducing costs.
But we must also consider the ethical implications. Data privacy is crucial, and we need strong regulations.
Finally, I believe that AI will create more jobs than it displaces if we invest in education.
Thank you for your attention!"""

@app.route("/demo-session", methods=['GET'])
def demo_session():
    """Get demo session with sample presentation analysis."""
    demo_analysis = {
        "word_count": 145,
        "wpm": 128,
        "filler_count": 4,
        "filler_details": {
            "uh": 2,
            "like": 1,
            "so": 1
        },
        "clarity": 82,
        "pace": "Good",
        "duration": 67.8,
        "script": DEMO_SCRIPT,
        "transcript": """Good morning, everyone! Today I want to talk about artificial intelligence and how it's transforming our world.
First, let me start with the basics. AI is uh, not a new concept, but recent advancements have made it incredibly powerful.
There are three key applications we should consider: healthcare, education, and business automation.
In healthcare, AI is helping doctors diagnose diseases like, faster and more accurately.
In education, personalized learning systems are adapting to each student's needs.
And in business, uh, automation is increasing efficiency and reducing costs.
But we must also consider the ethical implications. Data privacy is crucial, and we need strong regulations.
Finally, I believe that AI will create more jobs than it displaces so, if we invest in education.
Thank you for your attention!""",
        "suggestions": [
            "Great job! Your pace is excellent at 128 WPM - right in the optimal range!",
            "You used 4 filler words total. Try replacing them with strategic pauses for even better clarity.",
            "Your clarity score is 82/100 - very strong! Consider adding more emphasis on key statistics.",
            "Excellent presentation structure with clear transitions between topics.",
            "Pro tip: Slow down slightly when introducing new concepts to give audience time to absorb."
        ],
        "is_demo": True
    }
    return jsonify(demo_analysis)


@app.route("/demo-chat", methods=['POST'])
def demo_chat():
    """Get demo AI coach response (no API key required)."""
    data = request.get_json(silent=True) or {}
    user_message = data.get("message", "").lower().strip()
    
    # Predefined demo responses based on keywords
    demo_responses = {
        "filler": "Great question! Filler words like 'um' and 'uh' are natural, but try replacing them with pauses. When you feel one coming, take a breath instead - your audience will appreciate the moment to reflect on what you've said.",
        "pace": "Your pace is perfect at 128 WPM! The sweet spot is between 120-150 words per minute. If you need to slow down, add more descriptive details or pause after key points.",
        "clarity": "Your clarity score is 82/100 - excellent! To push it even higher, minimize filler words and add emphasis to important statistics. Your audience will find the content more memorable.",
        "confidence": "You're doing amazing! Your presentation flows well and your message is clear. The more you practice, the more confident you'll feel. Keep recording sessions to track your improvement!",
        "improve": "Here are your top improvement areas: 1) Reduce filler words by replacing them with pauses, 2) Add more varied intonation to emphasize key points, 3) Practice your transitions between topics.",
        "tips": "Best tips for better presentations: 1) Breathe naturally between sentences, 2) Practice beforehand, 3) Make eye contact with your audience, 4) Pace yourself between 120-150 WPM, 5) Use strategic pauses instead of filler words.",
        "default": "That's a great question about your speaking skills! Based on your demo session, you're doing well overall. Your pace is good and your clarity is strong. Keep practicing and you'll continue to improve!"
    }
    
    # Find matching response
    reply = demo_responses.get("default")
    for key, response in demo_responses.items():
        if key in user_message:
            reply = response
            break
    
    return jsonify({
        "reply": reply,
        "is_demo": True
    })


# ─── Run ──────────────────────────────────────────────────

def create_test_user():
    """Create a test user for demo purposes if none exists."""
    test_email = 'test@example.com'
    existing_user = User.query.filter_by(email=test_email).first()
    
    if not existing_user:
        test_user = User(
            email=test_email,
            username='testuser',
            full_name='Test User',
            email_verified=True
        )
        test_user.set_password('password123')
        db.session.add(test_user)
        db.session.commit()
        print("[✓] Test user created: test@example.com / password123")
    else:
        print("[✓] Test user already exists")


@app.before_request
def before_request():
    """Initialize database before first request."""
    if not hasattr(app, 'db_initialized'):
        with app.app_context():
            db.create_all()
            create_test_user()
            app.db_initialized = True

if __name__ == "__main__":
    debug_mode = os.getenv("FLASK_DEBUG", "true").lower() == "true"
    port = int(os.getenv("PORT", 5000))
    
    # Initialize database and create test user
    with app.app_context():
        db.create_all()
        create_test_user()
    
    app.run(debug=debug_mode, host="0.0.0.0", port=port)
